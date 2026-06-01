"""
extract_content.py
==================
Two-pass extraction pipeline for .pptx lecture files.

Pass 1: python-pptx extracts text boxes, titles, speaker notes, and exports
        all embedded images per slide. Flags slides as "text-light" when
        the combined extracted text is < 50 chars AND at least one image exists.

Pass 2: For every text-light slide, sends the exported image(s) to the
        Gemini Vision API (gemini-2.5-flash) for OCR / content description,
        and stores the result in the slide's `ocrContent` field.

Outputs
-------
- src/data/lectures/lec01.json … lec13.json
- public/images/lectures/lecXX/<slideNN>_img<N>.png
"""

import os
import re
import json
import time
import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from google import genai
from google.genai import types

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
BASE_DIR     = Path(__file__).resolve().parent.parent
LECTURES_DIR = BASE_DIR / "Lectures"
OUT_JSON_DIR = BASE_DIR / "src" / "data" / "lectures"
OUT_IMG_DIR  = BASE_DIR / "public" / "images" / "lectures"

# NOTE: Loaded from environment variables for push protection
GEMINI_API_KEY       = os.environ.get("GEMINI_API_KEY", "").strip()
GEMINI_MODEL         = "gemini-2.5-flash-lite"
TEXT_LIGHT_THRESHOLD = 50   # chars; below this the slide is flagged
RATE_LIMIT_SLEEP     = 7    # seconds between Gemini API calls (15 RPM free tier)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Gemini client setup (new google-genai SDK)
# ---------------------------------------------------------------------------
_client = genai.Client(api_key=GEMINI_API_KEY.strip())

# Validate the API key once at startup; disable OCR pass if invalid.
_GEMINI_AVAILABLE = False
try:
    _test_resp = _client.models.generate_content(
        model=GEMINI_MODEL, contents="ping"
    )
    _GEMINI_AVAILABLE = True
    log.info("Gemini API key validated — OCR pass will run.")
except Exception as _e:
    log.warning(
        "Gemini API key validation FAILED (%s). "
        "Pass 2 (OCR) will be SKIPPED. "
        "Please update GEMINI_API_KEY in this script and re-run.",
        _e,
    )


# ---------------------------------------------------------------------------
# Helper: parse filename → (lecture_number, week_number, title)
# ---------------------------------------------------------------------------
FNAME_RE = re.compile(
    r"Lec\s+(\d+)\s*-\s*W\s*(\d+)\s*-\s*(.+)\.pptx$", re.IGNORECASE
)

def parse_lecture_filename(fname: str):
    m = FNAME_RE.match(fname)
    if not m:
        log.warning("Cannot parse filename: %s", fname)
        return None, None, fname.replace(".pptx", "")
    return int(m.group(1)), int(m.group(2)), m.group(3).strip()


# ---------------------------------------------------------------------------
# Helper: extract text from a single pptx shape (recursive for groups)
# ---------------------------------------------------------------------------
def extract_shape_text(shape) -> str:
    parts = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            parts.append(extract_shape_text(s))
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = " ".join(run.text for run in para.runs if run.text.strip())
            if line.strip():
                parts.append(line.strip())
    return "\n".join(p for p in parts if p)


# ---------------------------------------------------------------------------
# Helper: export all images from a slide, return list of relative paths
# ---------------------------------------------------------------------------
def export_slide_images(slide, slide_dir: Path, slide_num: int) -> list:
    slide_dir.mkdir(parents=True, exist_ok=True)
    image_paths = []
    img_idx = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img_data  = shape.image.blob
                img_ext   = shape.image.ext or "png"
                img_fname = f"slide{slide_num:02d}_img{img_idx}.{img_ext}"
                img_path  = slide_dir / img_fname
                img_path.write_bytes(img_data)
                rel_path  = str(img_path.relative_to(BASE_DIR)).replace("\\", "/")
                image_paths.append(rel_path)
                img_idx += 1
            except Exception as exc:
                log.warning("  Could not export image on slide %d: %s", slide_num, exc)
    return image_paths


# ---------------------------------------------------------------------------
# Pass 1: Extract text + images for a single .pptx file
# ---------------------------------------------------------------------------
def pass1_extract(pptx_path: Path, lec_id: str):
    log.info("Pass 1 → %s", pptx_path.name)
    prs = Presentation(str(pptx_path))
    lec_img_dir = OUT_IMG_DIR / lec_id
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title   = ""
        content_lines = []
        notes_text    = ""

        # --- Extract slide title ---
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_title = slide.shapes.title.text_frame.text.strip()

        # --- Extract all text from all shapes ---
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            text = extract_shape_text(shape)
            if text:
                for line in text.split("\n"):
                    line = line.strip()
                    if line:
                        content_lines.append(line)

        # --- Extract speaker notes ---
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                raw_notes = notes_tf.text.strip()
                if raw_notes:
                    notes_text = raw_notes

        # --- Export images ---
        image_paths = export_slide_images(slide, lec_img_dir, idx)

        # --- Determine text-light status ---
        combined_text = slide_title + " " + " ".join(content_lines)
        is_text_light = (
            len(combined_text.strip()) < TEXT_LIGHT_THRESHOLD
            and len(image_paths) > 0
        )

        slides_data.append({
            "slideNumber":  idx,
            "title":        slide_title,
            "content":      content_lines,
            "ocrContent":   "",
            "notes":        notes_text,
            "images":       image_paths,
            "isTextLight":  is_text_light,
            "isScreenshot": False,
        })

    return slides_data


# ---------------------------------------------------------------------------
# Pass 2: OCR text-light slides via local Windows Media OCR Tool
# ---------------------------------------------------------------------------
def run_local_ocr(img_path: Path) -> str:
    """Run local OcrTool.exe to perform offline Windows OCR."""
    import subprocess
    tool_path = BASE_DIR / "scripts" / "OcrTool.exe"
    if not tool_path.exists():
        log.error("  OcrTool.exe not found at %s", tool_path)
        return "[OCR ERROR: OcrTool.exe not found]"
    
    try:
        res = subprocess.run(
            [str(tool_path), str(img_path)],
            capture_output=True,
            text=True,
            encoding="utf-8",
            check=True
        )
        return res.stdout.strip()
    except Exception as e:
        log.error("  Local OCR failed for %s: %s", img_path.name, e)
        stderr_msg = getattr(e, 'stderr', '') or ''
        return f"[OCR ERROR: {e} {stderr_msg}]"


def pass2_ocr(slides_data: list, lec_id: str):
    """Run local OCR on text-light slides. Returns enriched slides + stats."""
    text_light_slides = [s for s in slides_data if s["isTextLight"]]
    stats = {
        "total_slides":      len(slides_data),
        "text_light_slides": len(text_light_slides),
        "ocr_success":       0,
        "ocr_errors":        0,
    }

    if not text_light_slides:
        log.info("  No text-light slides — skipping OCR pass.")
        return slides_data, stats

    log.info("  %d text-light slide(s) found — starting local OCR …", len(text_light_slides))

    for slide in text_light_slides:
        slide_num = slide["slideNumber"]
        if not slide["images"]:
            continue

        first_img_rel = slide["images"][0]
        first_img_abs = BASE_DIR / first_img_rel.replace("/", os.sep)

        try:
            log.info("    OCR slide %d (%s) via local tool …", slide_num, first_img_abs.name)
            ocr_text = run_local_ocr(first_img_abs)
            slide["ocrContent"]   = ocr_text
            slide["isScreenshot"] = True
            stats["ocr_success"] += 1
        except Exception as exc:
            log.error("    OCR failed on slide %d: %s", slide_num, exc)
            slide["ocrContent"] = f"[OCR ERROR: {exc}]"
            stats["ocr_errors"] += 1

    return slides_data, stats



# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def process_lecture(pptx_path: Path):
    fname = pptx_path.name
    lec_num, week_num, title = parse_lecture_filename(fname)

    if lec_num is None:
        log.error("Skipping unrecognised file: %s", fname)
        return None

    lec_id = f"lec{lec_num:02d}"
    log.info("=" * 60)
    log.info("Processing %s  →  %s", fname, lec_id)

    slides_data           = pass1_extract(pptx_path, lec_id)
    slides_data, stats    = pass2_ocr(slides_data, lec_id)

    lecture_json = {
        "id":     lec_id,
        "number": lec_num,
        "week":   week_num,
        "title":  title,
        "slides": slides_data,
        "_stats": stats,
    }

    out_path = OUT_JSON_DIR / f"{lec_id}.json"
    out_path.write_text(
        json.dumps(lecture_json, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    log.info(
        "  Written → %s  (%d slides, %d text-light, %d OCR ok, %d err)",
        out_path.name,
        stats["total_slides"], stats["text_light_slides"],
        stats["ocr_success"],  stats["ocr_errors"],
    )
    return stats


def main():
    pptx_files = sorted(LECTURES_DIR.glob("*.pptx"))
    if not pptx_files:
        log.error("No .pptx files found in %s", LECTURES_DIR)
        return

    log.info("Found %d lecture file(s)", len(pptx_files))
    all_stats = []

    for pptx_path in pptx_files:
        result = process_lecture(pptx_path)
        if result:
            all_stats.append(result)

    log.info("=" * 60)
    log.info("DONE — global summary:")
    log.info("  Lectures processed : %d", len(all_stats))
    log.info("  Total slides       : %d", sum(s["total_slides"]      for s in all_stats))
    log.info("  Text-light slides  : %d", sum(s["text_light_slides"] for s in all_stats))
    log.info("  OCR success        : %d", sum(s["ocr_success"]       for s in all_stats))
    log.info("  OCR errors         : %d", sum(s["ocr_errors"]        for s in all_stats))


if __name__ == "__main__":
    main()
